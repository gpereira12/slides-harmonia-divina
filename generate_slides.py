import sys
import os
import re
import json
import requests
from html import unescape
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

load_dotenv()

# Default mass to fetch, override with argument: python generate_slides.py "Quaresma I"
MASS_NAME = sys.argv[1] if len(sys.argv) > 1 else "Quaresma II"
print(f"Gerando slides para: {MASS_NAME}...")

# 1. Fetching from Notion
token = os.getenv("NOTION_TOKEN")
database_id = os.getenv("NOTION_DATABASE_ID")
if not token or not database_id:
    raise ValueError("NOTION_TOKEN e NOTION_DATABASE_ID devem estar definidos no arquivo .env")
url = f"https://api.notion.com/v1/databases/{database_id}/query"
headers = {
    "Authorization": f"Bearer {token}",
    "Notion-Version": "2022-06-28",
    "Content-Type": "application/json"
}

all_results = []
has_more = True
next_cursor = None
while has_more:
    payload = {
        "filter": {
            "property": "Missa",
            "multi_select": {
                "contains": MASS_NAME
            }
        }
    }
    if next_cursor:
        payload["start_cursor"] = next_cursor
        
    response = requests.post(url, headers=headers, json=payload)
    if response.status_code == 200:
        data = response.json()
        results = data.get('results', [])
        all_results.extend(results)
        has_more = data.get('has_more', False)
        next_cursor = data.get('next_cursor', None)
    else:
        print(f"Error {response.status_code} fetching Notion: {response.text}")
        break

songs = []
for r in all_results:
    props = r.get("properties", {})
    title = props.get("Título", {}).get("title", [])
    title_text = title[0].get("plain_text", "") if title else "Untitled"
    momento = props.get("Momento", {}).get("multi_select", [])
    momento_texts = [m.get("name") for m in momento]
    letra = props.get("Letra Traduzida", {}).get("rich_text", [])
    letra_text = "".join(l.get("plain_text", "") for l in letra)
    
    songs.append({
        "title": title_text,
        "momento": momento_texts,
        "letra": letra_text,
        "letra_rich": letra
    })

# Default sorting: assuming Momento returns things like "01. Entrada"
songs.sort(key=lambda x: x["momento"][0] if x["momento"] else "99.")

# 2. Generating HTML
html_template = """<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Missa - Slides</title>

    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,400;0,600;1,400&family=Source+Serif+4:opsz,wght@8..60,400;8..60,600&display=swap" rel="stylesheet">

    <style>
        html, body {
            height: 100%;
            overflow-x: hidden;
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        html {
            scroll-snap-type: y mandatory;
            scroll-behavior: smooth;
        }

        :root {
            --bg-primary: #faf9f7;
            --text-primary: #1a1a1a;
            --text-secondary: #555555;
            --accent-crimson: #c41e3a;

            --font-display: 'Cormorant Garamond', serif;
            --font-body: 'Source Serif 4', serif;

            /* Greatly increased sizes for church projector legibility (base ~42px desktop) */
            --title-size: clamp(3.5rem, 8vw, 6.5rem);
            --h2-size: clamp(2.5rem, 6vw, 4.5rem);
            --body-size: clamp(2rem, 4vw, 3.5rem); /* ~42px to 56px */

            --slide-padding: clamp(2rem, 5vw, 5rem);
            --content-gap: clamp(1.5rem, 3vw, 2.5rem);
        }

        body {
            font-family: var(--font-body);
            background: var(--bg-primary);
            color: var(--text-primary);
            line-height: 1.35;
        }

        .slide {
            width: 100vw;
            height: 100vh;
            height: 100dvh;
            overflow: hidden;
            scroll-snap-align: start;
            display: flex;
            flex-direction: column;
            position: relative;
        }
        
        .slide.black-bg {
            background-color: #000;
        }

        .slide-content {
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            text-align: center;
            max-height: 100%;
            overflow: hidden;
            padding: var(--slide-padding);
            
            /* Max width constraints for precise readable blocks */
            width: 100%;
            max-width: 1024px;
            margin: 0 auto;
        }

        .title-slide {
            justify-content: center;
        }

        h1 {
            font-family: var(--font-display);
            font-size: var(--title-size);
            font-weight: 600;
            margin-bottom: 1rem;
            color: var(--text-primary);
        }

        .subtitle {
            font-family: var(--font-body);
            font-size: clamp(1.4rem, 2.5vw, 2rem);
            color: var(--accent-crimson);
            font-style: italic;
            letter-spacing: 0.05em;
            margin-top: 1rem;
            text-transform: uppercase;
        }

        .divider {
            width: 80px;
            height: 3px;
            background-color: var(--accent-crimson);
            margin: 1.5rem auto;
            opacity: 0.7;
        }

        .lyric-stanza {
            font-size: var(--body-size);
            line-height: 1.35;
            margin-bottom: var(--content-gap);
            width: 100%;
            white-space: pre-wrap;
        }

        .lyric-stanza:last-child {
            margin-bottom: 0;
        }

        @media (max-width: 768px) {
            .slide-content {
                max-width: 600px;
            }
        }
        @media (max-height: 700px) {
            :root {
                --slide-padding: clamp(1rem, 3vw, 2rem);
                --content-gap: clamp(0.5rem, 2vw, 1.5rem);
            }
        }
        @media (prefers-reduced-motion: reduce) {
            *, *::before, *::after {
                animation-duration: 0.01ms !important;
                transition-duration: 0.2s !important;
            }
            html { scroll-behavior: auto; }
        }

        .reveal {
            opacity: 0;
            transform: translateY(20px);
            transition: opacity 0.8s ease, transform 0.8s ease;
        }
        .slide.visible .reveal {
            opacity: 1;
            transform: translateY(0);
        }
        .delay-1 { transition-delay: 0.1s; }
        .delay-2 { transition-delay: 0.2s; }
    </style>
<body>
    <div id="presentation">
"""

html_body = f'''
        <!-- Main Title Slide -->
        <section class="slide title-slide">
            <div class="slide-content">
                <h1 class="reveal">{MASS_NAME}</h1>
                <div class="divider reveal delay-1"></div>
            </div>
        </section>
'''
slide_count = 1
slides_data = [{"type": "title", "text": MASS_NAME, "bold": False}]

for song in songs:
    momento = song.get("momento", [""])[0]
    title = song.get("title", "")
    letra = song.get("letra", "").strip()

    if not letra:
        continue # Ignora musicas sem letra traduzida
        
    if "fora da liturgia" in momento.lower():
        continue # Ignora cantos fora da liturgia

    if "aspersão" in momento.lower() or "aspersao" in momento.lower():
        continue # Ignora aspersão

    # Clean momento (e.g. "01. Entrada" -> "Entrada", "5.2 Cinzas" -> "Cinzas")
    momento_clean = re.sub(r'^\d+[\.\-]?\d*[\s\.\-]*', '', momento).strip()
        
    if title == "Untitled" or not title:
        display_title = momento_clean
    else:
        display_title = title

    # ==== Titulo da Musica ====
    # Slide preto divisorio antes
    html_body += f'''
        <section class="slide black-bg" id="slide-{slide_count}"></section>
    '''
    slides_data.append({"type": "black", "text": "", "bold": False})
    slide_count += 1
    
    if display_title.lower() == momento_clean.lower():
        # Evita duplicar a palavra quando o título é o mesmo que o momento
        html_body += f'''
            <section class="slide title-slide" id="slide-{slide_count}">
                <div class="slide-content">
                    <h1 class="reveal">{display_title}</h1>
                </div>
            </section>
        '''
        slides_data.append({"type": "song_title", "text": display_title, "subtitle": None, "bold": False})
    else:
        html_body += f'''
            <section class="slide title-slide" id="slide-{slide_count}">
                <div class="slide-content">
                    <p class="subtitle reveal">{momento_clean}</p>
                    <div class="divider reveal delay-1"></div>
                    <h1 class="reveal delay-2">{display_title}</h1>
                </div>
            </section>
        '''
        slides_data.append({"type": "song_title", "text": display_title, "subtitle": momento_clean, "bold": False})
    slide_count += 1

    letra_rich = song.get("letra_rich", [])
    has_bold = any(seg.get("annotations", {}).get("bold", False) for seg in letra_rich)

    # ==== REGRAS DE EXCEÇÃO GLOBAIS ====
    
    # 1. Ato Penitencial (Fórmula 2/3)
    if "Ato Penitencial" in momento_clean:
        if "fórmula 2" in display_title.lower() or "fórmula 2" in letra.lower():
            # Formula 2: Agrupa normal e negrito no mesmo slide (Solo e Resposta do povo)
            current_slide_content = []
            for seg in letra_rich:
                text = seg.get("plain_text", "").strip()
                if not text: continue
                is_bold = seg.get("annotations", {}).get("bold", False)
                text_html = text.replace('\\n', '<br>')
                
                if is_bold:
                    current_slide_content.append(f"<strong>{text_html}</strong>")
                    final_html = "<br>".join(current_slide_content)
                    html_body += f'''
                    <section class="slide" id="slide-{slide_count}">
                        <div class="slide-content">
                            <p class="lyric-stanza reveal delay-1">{final_html}</p>
                        </div>
                    </section>
                    '''
                    plain_text = ' / '.join(c.replace('<strong>', '').replace('</strong>', '') for c in current_slide_content)
                    slides_data.append({"type": "lyric", "text": plain_text, "bold": True})
                    slide_count += 1
                    current_slide_content = [] # Reset
                else:
                    current_slide_content.append(text_html)
                    
            if current_slide_content:
                plain_text = ' / '.join(current_slide_content)
                final_html = "<br>".join(current_slide_content)
                html_body += f'''
                <section class="slide" id="slide-{slide_count}">
                    <div class="slide-content">
                        <p class="lyric-stanza reveal delay-1">{final_html}</p>
                    </div>
                </section>
                '''
                slides_data.append({"type": "lyric", "text": plain_text, "bold": False})
                slide_count += 1
                
        else:
            # Formula 3 (e 1 se houver): Renderização SEQUENCIAL exata do que está no Notion
            for seg in letra_rich:
                text = seg.get("plain_text", "").strip()
                if not text: continue
                is_bold = seg.get("annotations", {}).get("bold", False)
                text_html = text.replace('\n', '<br>')
                
                if is_bold:
                    html_body += f'''
                    <section class="slide" id="slide-{slide_count}">
                        <div class="slide-content">
                            <p class="lyric-stanza reveal delay-1"><strong>{text_html}</strong></p>
                        </div>
                    </section>
                    '''
                    slides_data.append({"type": "lyric", "text": text, "bold": True})
                    slide_count += 1
                else:
                    stanzas = [p.strip() for p in text.split('\n\n') if p.strip()]
                    for s in stanzas:
                        s_html = s.replace('\n', '<br>')
                        html_body += f'''
                        <section class="slide" id="slide-{slide_count}">
                            <div class="slide-content">
                                <p class="lyric-stanza reveal delay-1">{s_html}</p>
                            </div>
                        </section>
                        '''
                        slides_data.append({"type": "lyric", "text": s, "bold": False})
                        slide_count += 1

    # 2. Aclamação ao Evangelho (ou Aleluia)
    elif "Aclamação" in momento_clean or "Aleluia" in momento_clean:
        for seg in letra_rich:
            is_bold = seg.get("annotations", {}).get("bold", False)
            if is_bold:
                text = seg.get("plain_text", "").strip()
                if not text: continue
                text_html = text.replace('\n', '<br>')
                html_body += f'''
                <section class="slide" id="slide-{slide_count}">
                    <div class="slide-content">
                        <p class="lyric-stanza reveal delay-1"><strong>{text_html}</strong></p>
                    </div>
                </section>
                '''
                slides_data.append({"type": "lyric", "text": text, "bold": True})
                slide_count += 1

    # 3. Santo e Cordeiro
    elif "Santo" in momento_clean or "Cordeiro" in momento_clean:
        for seg in letra_rich:
            text = seg.get("plain_text", "").strip()
            if not text: continue
            is_bold = seg.get("annotations", {}).get("bold", False)
            text_html = text.replace('\n', '<br>')
            
            if is_bold:
                html_body += f'''
                <section class="slide" id="slide-{slide_count}">
                    <div class="slide-content">
                        <p class="lyric-stanza reveal delay-1"><strong>{text_html}</strong></p>
                    </div>
                </section>
                '''
                slides_data.append({"type": "lyric", "text": text, "bold": True})
                slide_count += 1
            else:
                stanzas = [p.strip() for p in text.split('\n\n') if p.strip()]
                for s in stanzas:
                    s_html = s.replace('\n', '<br>')
                    html_body += f'''
                    <section class="slide" id="slide-{slide_count}">
                        <div class="slide-content">
                            <p class="lyric-stanza reveal delay-1">{s_html}</p>
                        </div>
                    </section>
                    '''
                    slides_data.append({"type": "lyric", "text": s, "bold": False})
                    slide_count += 1

    # 4. REGRA PADRÃO (Entrada, Ofertório, Comunhão, Final...)
    elif has_bold:
        chorus = ""
        raw_stanzas = []
        
        for seg in letra_rich:
            text = seg.get("plain_text", "").strip()
            if not text: continue
            
            is_bold = seg.get("annotations", {}).get("bold", False)
            if is_bold and not chorus:
                chorus = text
            elif is_bold and chorus:
                raw_stanzas.append(f"<strong>{text}</strong>")
            else:
                raw_stanzas.append(text)
                
        # Parse stanzas
        stanzas = []
        for text in raw_stanzas:
            parts = [p.strip() for p in text.split('\n\n') if p.strip()]
            for part in parts:
                lines = part.split('\n')
                current = []
                for line in lines:
                    line = line.strip()
                    if not line: continue
                    if re.match(r'^\\d+\\.', line) and current:
                        stanzas.append('<br>'.join(current))
                        current = [line]
                    else:
                        current.append(line)
                if current:
                    stanzas.append('<br>'.join(current))
                    
        # Render alternating chunks
        if chorus:
            chorus_html = chorus.replace('\n', '<br>')
            html_body += f'''
        <section class="slide" id="slide-{slide_count}">
            <div class="slide-content">
                <p class="lyric-stanza reveal delay-1"><strong>{chorus_html}</strong></p>
            </div>
        </section>
            '''
            slides_data.append({"type": "lyric", "text": chorus, "bold": True})
            slide_count += 1
            
            for stanza in stanzas:
                html_body += f'''
        <section class="slide" id="slide-{slide_count}">
            <div class="slide-content">
                <p class="lyric-stanza reveal delay-1">{stanza}</p>
            </div>
        </section>
                '''
                stanza_plain = re.sub(r'<[^>]+>', '', stanza)
                slides_data.append({"type": "lyric", "text": stanza_plain, "bold": False})
                slide_count += 1
                
                html_body += f'''
        <section class="slide" id="slide-{slide_count}">
            <div class="slide-content">
                <p class="lyric-stanza reveal delay-1"><strong>{chorus_html}</strong></p>
            </div>
        </section>
                '''
                slides_data.append({"type": "lyric", "text": chorus, "bold": True})
                slide_count += 1
                
    else:
        stanzas = []
        parts = [p.strip() for p in letra.split('\n\n') if p.strip()]
        for part in parts:
            lines = part.split('\n')
            current = []
            for line in lines:
                line = line.strip()
                if not line: continue
                if re.match(r'^\\d+\\.', line) and current:
                    stanzas.append('<br>'.join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                stanzas.append('<br>'.join(current))
        
        for p_html in stanzas:
            html_body += f'''
        <section class="slide" id="slide-{slide_count}">
            <div class="slide-content">
                <p class="lyric-stanza reveal delay-1">{p_html}</p>
            </div>
        </section>
            '''
            plain_text = re.sub(r'<[^>]+>', '', p_html)
            slides_data.append({"type": "lyric", "text": plain_text, "bold": False})
            slide_count += 1

html_footer = """
    </div>

    <script>
        document.addEventListener("DOMContentLoaded", () => {
            const slides = document.querySelectorAll('.slide');
            
            const observer = new IntersectionObserver((entries) => {
                entries.forEach(entry => {
                    if (entry.isIntersecting) {
                        entry.target.classList.add('visible');
                    }
                });
            }, { root: null, rootMargin: '0px', threshold: 0.5 });

            slides.forEach(slide => observer.observe(slide));

            document.addEventListener('keydown', (e) => {
                let currentIndex = 0;
                let maxVisibleArea = 0;
                
                slides.forEach((slide, index) => {
                    const rect = slide.getBoundingClientRect();
                    const visibleHeight = Math.min(rect.bottom, window.innerHeight) - Math.max(rect.top, 0);
                    if (visibleHeight > maxVisibleArea) {
                        maxVisibleArea = visibleHeight;
                        currentIndex = index;
                    }
                });

                if (e.key === 'ArrowDown' || e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown') {
                    if (currentIndex < slides.length - 1) {
                        e.preventDefault();
                        slides[currentIndex + 1].scrollIntoView({ behavior: 'smooth' });
                    }
                } else if (e.key === 'ArrowUp' || e.key === 'ArrowLeft' || e.key === 'PageUp') {
                    if (currentIndex > 0) {
                        e.preventDefault();
                        slides[currentIndex - 1].scrollIntoView({ behavior: 'smooth' });
                    }
                }
            });
            
            document.body.addEventListener('click', (e) => {
                if (e.target.classList.contains('dot')) return;
                let currentIndex = 0;
                let maxVisibleArea = 0;
                slides.forEach((slide, index) => {
                    const rect = slide.getBoundingClientRect();
                    const visibleHeight = Math.min(rect.bottom, window.innerHeight) - Math.max(rect.top, 0);
                    if (visibleHeight > maxVisibleArea) {
                        maxVisibleArea = visibleHeight;
                        currentIndex = index;
                    }
                });
                if (currentIndex < slides.length - 1) {
                    slides[currentIndex + 1].scrollIntoView({ behavior: 'smooth' });
                }
            });
        });
    </script>
</body>
</html>
"""

os.makedirs("/Users/gabrielpereira/Desktop/Slides/missas/Quaresma", exist_ok=True)
output_file = f"/Users/gabrielpereira/Desktop/Slides/missas/Quaresma/{MASS_NAME.replace(' ', '_')}.html"
with open(output_file, "w") as f:
    f.write(html_template + html_body + html_footer)

print(f"Generated {slide_count} HTML slides for {MASS_NAME}.")

# ============================================================
# PPTX EXPORT
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette - cream/warm light theme
BG_CREAM  = RGBColor(0xFA, 0xF9, 0xF7)  # fundo creme
BG_BLACK  = RGBColor(0x00, 0x00, 0x00)  # slide preto de transicao
CLR_DARK  = RGBColor(0x1A, 0x1A, 0x1A)  # texto principal
CLR_RED   = RGBColor(0xC4, 0x1E, 0x3A)  # acento vermelho
CLR_GREY  = RGBColor(0x44, 0x44, 0x44)  # texto secundario (estrofes)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_centered_textbox(slide, text, font_size, color, bold=False, italic=False):
    """Full-slide textbox, text centered vertically and horizontally."""
    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.8)
    txBox = slide.shapes.add_textbox(M, M, W - 2*M, H - 2*M)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox

def add_subtitle_and_title(slide, title, subtitle=None):
    """Song title slide with optional subtitle above."""
    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.8)
    if subtitle:
        sub_box = slide.shapes.add_textbox(M, H * 0.28, W - 2*M, H * 0.15)
        tf = sub_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle.upper()
        run.font.size = Pt(20)
        run.font.color.rgb = CLR_RED
        run.font.italic = True

    title_box = slide.shapes.add_textbox(M, H * 0.42 if subtitle else H * 0.3, W - 2*M, H * 0.35)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(52)
    run.font.color.rgb = CLR_DARK
    run.font.bold = True

blank_layout = prs.slide_layouts[6]  # blank

for sd in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    stype = sd["type"]

    if stype == "black":
        add_bg(slide, BG_BLACK)

    elif stype == "title":
        add_bg(slide, BG_CREAM)
        add_subtitle_and_title(slide, sd["text"])

    elif stype == "song_title":
        add_bg(slide, BG_CREAM)
        add_subtitle_and_title(slide, sd["text"], sd.get("subtitle"))

    elif stype == "lyric":
        add_bg(slide, BG_CREAM)
        text_color = CLR_DARK if sd.get("bold") else CLR_GREY
        font_sz = 36 if sd.get("bold") else 32
        add_centered_textbox(slide, sd["text"], font_sz, text_color, bold=sd.get("bold", False))

pptx_file = output_file.replace(".html", ".pptx")
prs.save(pptx_file)
print(f"Generated PPTX: {pptx_file}")
print(f"Done! Total: {slide_count} slides.")
